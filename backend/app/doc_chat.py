import io
import fitz
from docx import Document as DocxDocument
from pptx import Presentation
from backend.app.groq_client import chat as groq_chat

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_ollama import OllamaEmbeddings
from langchain_core.documents import Document
from datetime import datetime

embeddings = OllamaEmbeddings(model="nomic-embed-text")

SUPPORTED_TYPES = {
    ".pdf": "PDF Document",
    ".docx": "Word Document",
    ".doc": "Word Document",
    ".pptx": "PowerPoint Presentation",
    ".ppt": "PowerPoint Presentation",
    ".txt": "Text File",
}

def extract_text(contents: bytes, filename: str) -> dict:
    """Extract text from any supported document type."""
    ext = "." + filename.lower().rsplit(".", 1)[-1]

    if ext == ".pdf":
        pdf_doc = fitz.open(stream=contents, filetype="pdf")
        pages = []
        for i, page in enumerate(pdf_doc):
            text = page.get_text()
            if text.strip():
                pages.append({"page": i + 1, "text": text})
        pdf_doc.close()
        return {
            "type": "PDF Document",
            "pages": pages,
            "full_text": "\n".join([p["text"] for p in pages]),
            "page_count": len(pages),
        }

    elif ext in [".docx", ".doc"]:
        doc = DocxDocument(io.BytesIO(contents))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                if row_text:
                    tables_text.append(row_text)
        full_text = "\n".join(paragraphs + tables_text)
        return {
            "type": "Word Document",
            "pages": [{"page": 1, "text": full_text}],
            "full_text": full_text,
            "page_count": 1,
            "paragraphs": len(paragraphs),
            "tables": len(doc.tables),
        }

    elif ext in [".pptx", ".ppt"]:
        prs = Presentation(io.BytesIO(contents))
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if shape.shape_type == 13:
                        continue
                    if "title" in shape.name.lower():
                        slide_title = shape.text.strip()
                    else:
                        slide_text.append(shape.text.strip())
            if slide_title or slide_text:
                slides.append({
                    "page": i + 1,
                    "title": slide_title,
                    "text": f"Slide {i+1}: {slide_title}\n" + "\n".join(slide_text),
                })
        return {
            "type": "PowerPoint Presentation",
            "pages": slides,
            "full_text": "\n\n".join([s["text"] for s in slides]),
            "page_count": len(prs.slides),
            "slide_count": len(slides),
        }

    elif ext == ".txt":
        text = contents.decode("utf-8", errors="ignore")
        return {
            "type": "Text File",
            "pages": [{"page": 1, "text": text}],
            "full_text": text,
            "page_count": 1,
        }

    else:
        raise ValueError(f"Unsupported file type: {ext}")

def generate_summary(doc_data: dict, filename: str) -> str:
    """Generate AI summary of the document."""
    sample = doc_data["full_text"][:3000]
    doc_type = doc_data["type"]

    prompt = f"""You are an expert document analyst. Analyze this {doc_type} and provide:

1. DOCUMENT TYPE & PURPOSE (1-2 sentences)
2. MAIN TOPICS (3-5 bullet points)
3. KEY INFORMATION (important facts, data, conclusions)
4. DOCUMENT STRUCTURE (how it's organized)

Document: {filename}
Content:
{sample}

Be concise and professional."""

    return groq_chat(prompt)

def ingest_document(contents: bytes, filename: str, file_id: str) -> dict:
    """Extract text and index into ChromaDB."""
    doc_data = extract_text(contents, filename)

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500, chunk_overlap=50
    )
    chunks = splitter.split_text(doc_data["full_text"])

    docs = [
        Document(
            page_content=chunk,
            metadata={"source": filename, "file_id": file_id, "chunk": i}
        )
        for i, chunk in enumerate(chunks)
    ]

    collection_name = f"doc_{file_id[:20].replace('-', '_')}"
    vectorstore = Chroma(
        collection_name=collection_name,
        embedding_function=embeddings,
        persist_directory="./chroma_db"
    )
    vectorstore.add_documents(docs)

    summary = generate_summary(doc_data, filename)

    return {
        "filename": filename,
        "doc_type": doc_data["type"],
        "page_count": doc_data["page_count"],
        "chunks_indexed": len(chunks),
        "summary": summary,
        "metadata": {k: v for k, v in doc_data.items() if k not in ["full_text", "pages"]},
    }

def query_document(question: str, file_id: str) -> dict:
    """Query document using RAG."""
    collection_name = f"doc_{file_id[:20].replace('-', '_')}"

    vectorstore = Chroma(
        collection_name=collection_name,
        embedding_function=embeddings,
        persist_directory="./chroma_db"
    )

    retriever = vectorstore.as_retriever(search_kwargs={"k": 4})
    relevant_docs = retriever.invoke(question)
    context = "\n\n".join([doc.page_content for doc in relevant_docs])

    prompt = f"""You are an expert document assistant. Answer the question based ONLY on the document content below.

DOCUMENT CONTENT:
{context}

QUESTION: {question}

If the answer is not in the document, say "This information is not found in the document."
Provide a clear, accurate, and helpful answer."""

    answer = groq_chat(prompt)

    return {
        "question": question,
        "answer": answer,
        "context_chunks": len(relevant_docs),
    }
